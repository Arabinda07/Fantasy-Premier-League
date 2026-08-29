import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_fpl_pptx():
    prs = Presentation()
    # 4:5 Vertical Portrait Ratio (1080 x 1350 px equivalent in inches: 8.0 in x 10.0 in)
    prs.slide_width = Inches(8.0)
    prs.slide_height = Inches(10.0)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Tokens
    C_BLUE = RGBColor(6, 106, 254)         # #066AFE
    C_COBALT = RGBColor(46, 53, 238)       # #2E35EE
    C_NAVY_DARK = RGBColor(8, 15, 38)      # #080F26
    C_OBSIDIAN = RGBColor(8, 12, 26)       # #080C1A
    C_INDIGO = RGBColor(10, 19, 51)        # #0A1333
    C_SUNSET_DARK = RGBColor(28, 15, 8)    # #1C0F08
    C_TURF = RGBColor(14, 41, 24)          # #0E2918
    C_MINT = RGBColor(0, 163, 137)         # #00A389
    C_TANGERINE = RGBColor(255, 122, 0)    # #FF7A00
    C_ORCHID = RGBColor(224, 64, 251)      # #E040FB
    C_AMBER = RGBColor(255, 176, 32)       # #FFB020
    C_WHITE = RGBColor(255, 255, 255)      # #FFFFFF
    C_SLATE = RGBColor(148, 163, 184)      # #94A3B8
    C_LAVENDER = RGBColor(196, 181, 253)   # #C4B5FD
    C_CARD_BG = RGBColor(20, 27, 48)       # #141B30

    assets_dir = os.path.abspath("carousels/fpl-algorithm-experiment/assets")
    screenshot_path = os.path.join(assets_dir, "fpl-lineup-screenshot.png")

    def add_header_footer(slide, slide_num, total_slides=6):
        # Header
        h_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7.0), Inches(0.4))
        tf = h_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"0{slide_num} / 0{total_slides}   •   ARABINDA SAHA"
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_SLATE

        # Footer
        f_box = slide.shapes.add_textbox(Inches(0.5), Inches(9.3), Inches(7.0), Inches(0.4))
        tf = f_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "@robin0607saha   •   Swipe →" if slide_num < 6 else "@robin0607saha   •   Completed ✓"
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_MINT if slide_num == 6 else C_SLATE

    # ==========================================
    # SLIDE 1: THE HOOK / COVER (Text Only)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.0), Inches(10.0))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_NAVY_DARK
    bg1.line.color.rgb = C_NAVY_DARK

    add_header_footer(s1, 1)

    # Tag Badge
    badge1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.0), Inches(3.2), Inches(0.4))
    badge1.fill.solid()
    badge1.fill.fore_color.rgb = RGBColor(25, 45, 95)
    badge1.line.color.rgb = C_BLUE
    p = badge1.text_frame.paragraphs[0]
    p.text = "SYSTEM EXPERIMENT // 2026-27"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Main Headline
    t_box1 = s1.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(7.0), Inches(2.8))
    tf1 = t_box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "I MISSED GAMEWEEK 1.\nSO I'M LETTING AN ALGORITHM RUN MY ENTIRE FPL SEASON FROM 0 POINTS."
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "37 Gameweeks. £100.0M Budget. Zero Human Emotion."
    p1_sub.font.size = Pt(17)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = C_AMBER
    p1_sub.space_before = Pt(12)

    # 3 Column Pillars
    col_w = Inches(2.2)
    gap = Inches(0.2)
    pillars = [
        ("PILLAR 01", "Match Data Ingestion", "Pulls underlying xG, xA, shot locations, and rotation risk factors."),
        ("PILLAR 02", "11 Probability Models", "Deconstructs points into discrete Poisson-adjusted match events."),
        ("PILLAR 03", "Integer Linear Solver", "Optimizes the 15-man squad across a 5-week rolling horizon.")
    ]

    for i, (tag, title, desc) in enumerate(pillars):
        left_pos = Inches(0.5) + i * (col_w + gap)
        card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(4.8), col_w, Inches(3.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = RGBColor(40, 55, 95)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.15)
        ctf.margin_right = Inches(0.15)
        ctf.margin_top = Inches(0.2)
        
        cp1 = ctf.paragraphs[0]
        cp1.text = tag
        cp1.font.size = Pt(11)
        cp1.font.bold = True
        cp1.font.color.rgb = C_MINT

        cp2 = ctf.add_paragraph()
        cp2.text = title
        cp2.font.size = Pt(15)
        cp2.font.bold = True
        cp2.font.color.rgb = C_WHITE
        cp2.space_before = Pt(8)

        cp3 = ctf.add_paragraph()
        cp3.text = desc
        cp3.font.size = Pt(12)
        cp3.font.color.rgb = C_SLATE
        cp3.space_before = Pt(8)

    # Takeaway Banner
    ban1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(8.4), Inches(7.0), Inches(0.7))
    ban1.fill.solid()
    ban1.fill.fore_color.rgb = RGBColor(12, 20, 45)
    ban1.line.color.rgb = C_BLUE
    bp = ban1.text_frame.paragraphs[0]
    bp.text = "Can pure mathematics beat 11 million human managers over 37 weeks? →"
    bp.font.size = Pt(13)
    bp.font.bold = True
    bp.font.color.rgb = C_WHITE
    bp.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: THE EMOTION TRAP (Text Only)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    bg2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.0), Inches(10.0))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = C_OBSIDIAN
    bg2.line.color.rgb = C_OBSIDIAN

    add_header_footer(s2, 2)

    # Title
    t_box2 = s2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(7.0), Inches(1.2))
    tf2 = t_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "THE FRIDAY 1:00 AM TRAP"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE

    p2_sub = tf2.add_paragraph()
    p2_sub.text = "Why rational football fans make completely irrational fantasy choices."
    p2_sub.font.size = Pt(15)
    p2_sub.font.color.rgb = C_SLATE
    p2_sub.space_before = Pt(4)

    # 2 Column Cards (Human vs Model)
    col_w2 = Inches(3.35)
    gap2 = Inches(0.3)

    # Left: Human
    card_h = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.4), col_w2, Inches(5.6))
    card_h.fill.solid()
    card_h.fill.fore_color.rgb = RGBColor(28, 18, 12)
    card_h.line.color.rgb = C_TANGERINE
    
    htf = card_h.text_frame
    htf.word_wrap = True
    htf.margin_left = Inches(0.2)
    htf.margin_top = Inches(0.25)
    
    hp = htf.paragraphs[0]
    hp.text = "THE 5-DAY HUMAN LOOP"
    hp.font.size = Pt(13)
    hp.font.bold = True
    hp.font.color.rgb = C_TANGERINE

    steps_h = [
        ("MON - THU: Rational Analysis", "Study underlying numbers, watch press conferences, and build a logical plan."),
        ("FRIDAY 1:00 AM: The Knee-Jerk", "Throw the plan away. Captain a differential purely on a 10-second gut feeling."),
        ("SUNDAY 94TH MIN: The Heartbreak", "Clean sheet wiped out by a deflected cross. Rage-quit until next weekend.")
    ]
    for st, sd in steps_h:
        p_t = htf.add_paragraph()
        p_t.text = st
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE
        p_t.space_before = Pt(16)
        
        p_d = htf.add_paragraph()
        p_d.text = sd
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = C_SLATE
        p_d.space_before = Pt(3)

    # Right: Model
    card_m = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5) + col_w2 + gap2, Inches(2.4), col_w2, Inches(5.6))
    card_m.fill.solid()
    card_m.fill.fore_color.rgb = RGBColor(10, 30, 24)
    card_m.line.color.rgb = C_MINT
    
    mtf = card_m.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.2)
    mtf.margin_top = Inches(0.25)
    
    mp = mtf.paragraphs[0]
    mp.text = "THE MATHEMATICAL REMEDY"
    mp.font.size = Pt(13)
    mp.font.bold = True
    mp.font.color.rgb = C_MINT

    steps_m = [
        ("STEP 01: Objective Baseline Priors", "Evaluates shot quality and venue dynamics without human narrative bias."),
        ("STEP 02: Bayesian Shrinkage", "Blends 6-week short form with multi-season baselines to stop cameo overreactions."),
        ("STEP 03: 5-Week Horizon Solver", "Optimizes for rolling multi-week expectation rather than single-week panic.")
    ]
    for st, sd in steps_m:
        p_t = mtf.add_paragraph()
        p_t.text = st
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE
        p_t.space_before = Pt(16)
        
        p_d = mtf.add_paragraph()
        p_d.text = sd
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = C_SLATE
        p_d.space_before = Pt(3)

    # Takeaway Banner
    ban2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(8.3), Inches(7.0), Inches(0.7))
    ban2.fill.solid()
    ban2.fill.fore_color.rgb = RGBColor(12, 28, 22)
    ban2.line.color.rgb = C_MINT
    bp2 = ban2.text_frame.paragraphs[0]
    bp2.text = "Emotion maximizes variance. Mathematical solvers optimize expectation."
    bp2.font.size = Pt(13)
    bp2.font.bold = True
    bp2.font.color.rgb = C_WHITE
    bp2.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 3: 11 POINT PROBABILITIES (Tactical Diagram)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    bg3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.0), Inches(10.0))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = C_INDIGO
    bg3.line.color.rgb = C_INDIGO

    add_header_footer(s3, 3)

    t_box3 = s3.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(7.0), Inches(1.2))
    tf3 = t_box3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "DECONSTRUCTING POINTS INTO MATH"
    p3.font.size = Pt(26)
    p3.font.bold = True
    p3.font.color.rgb = C_WHITE

    p3_sub = tf3.add_paragraph()
    p3_sub.text = "The engine does not guess raw points. It models 11 distinct event probabilities."
    p3_sub.font.size = Pt(14)
    p3_sub.font.color.rgb = C_SLATE
    p3_sub.space_before = Pt(4)

    # Left: 3 Stacked Cards
    card_w3 = Inches(3.4)
    probs = [
        ("01. ATTACKING SUITE", C_TANGERINE, "xG90 (Goal threat by pitch zone)\nxA90 (Open play assist rate)\nSet Pieces (Penalties & corners)"),
        ("02. DEFENSE ENGINE", C_BLUE, "Poisson Clean Sheet Probability\nGoals Conceded Step Deductions\nGoalkeeper Save Volume Scaling"),
        ("03. REALITY LAYER", C_MINT, "Lineup Starting Probability\nDynamic Bonus Point (BPS) Capture\nYellow & Red Card Discipline")
    ]
    for idx, (head, col, body) in enumerate(probs):
        top_y = Inches(2.4) + idx * Inches(1.9)
        c_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top_y, card_w3, Inches(1.7))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = RGBColor(16, 26, 56)
        c_box.line.color.rgb = col
        
        ctf = c_box.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.18)
        ctf.margin_top = Inches(0.15)
        
        cp1 = ctf.paragraphs[0]
        cp1.text = head
        cp1.font.size = Pt(12)
        cp1.font.bold = True
        cp1.font.color.rgb = col
        
        cp2 = ctf.add_paragraph()
        cp2.text = body
        cp2.font.size = Pt(11.5)
        cp2.font.color.rgb = C_WHITE
        cp2.space_before = Pt(6)

    # Right: Tactical Pitch Penalty Box Diagram
    diag_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(2.4), Inches(3.4), Inches(5.5))
    diag_box.fill.solid()
    diag_box.fill.fore_color.rgb = C_TURF
    diag_box.line.color.rgb = RGBColor(40, 160, 90)
    
    dtf = diag_box.text_frame
    dtf.word_wrap = True
    dtf.margin_left = Inches(0.2)
    dtf.margin_top = Inches(0.25)
    
    dp1 = dtf.paragraphs[0]
    dp1.text = "TACTICAL PENALTY BOX MODEL"
    dp1.font.size = Pt(13)
    dp1.font.bold = True
    dp1.font.color.rgb = C_WHITE
    
    items_d = [
        ("Zone 1: Danger Box (xG 0.65)", "Central 6-yard box probability cluster."),
        ("Zone 2: Wide Assist Channels", "High-frequency cutback passing lanes."),
        ("Zone 3: Poisson Defensive Shield", "Shot suppression & clean sheet odds."),
        ("Zone 4: BPS Dynamic Capture", "Bonus distribution based on match dominance.")
    ]
    for dt, dd in items_d:
        p_t = dtf.add_paragraph()
        p_t.text = dt
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_AMBER
        p_t.space_before = Pt(14)
        
        p_d = dtf.add_paragraph()
        p_d.text = dd
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_SLATE
        p_d.space_before = Pt(2)

    # Banner
    ban3 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(8.3), Inches(7.0), Inches(0.7))
    ban3.fill.solid()
    ban3.fill.fore_color.rgb = RGBColor(16, 26, 56)
    ban3.line.color.rgb = C_BLUE
    bp3 = ban3.text_frame.paragraphs[0]
    bp3.text = "Granular micro-events compound into resilient season-long predictions."
    bp3.font.size = Pt(13)
    bp3.font.bold = True
    bp3.font.color.rgb = C_WHITE
    bp3.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 4: THE REAL SQUAD (Screenshot Diagram)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    bg4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.0), Inches(10.0))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = RGBColor(17, 5, 44)
    bg4.line.color.rgb = RGBColor(17, 5, 44)

    add_header_footer(s4, 4)

    t_box4 = s4.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(7.0), Inches(1.2))
    tf4 = t_box4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "THE ALGORITHM'S STARTING LINEUP"
    p4.font.size = Pt(26)
    p4.font.bold = True
    p4.font.color.rgb = C_WHITE

    p4_sub = tf4.add_paragraph()
    p4_sub.text = "The mathematically optimal 15-player squad within £100.0M."
    p4_sub.font.size = Pt(14)
    p4_sub.font.color.rgb = C_SLATE
    p4_sub.space_before = Pt(4)

    # Left: Rationale Cards
    card_w4 = Inches(3.2)
    squad_cards = [
        ("CAPTAINCY LOCK", C_AMBER, "Bruno Fernandes (C)", "Highest simulated 5-week expected yield vs Ipswich (H)."),
        ("MIDFIELD POWER", C_MINT, "Palmer • Mbeumo • Szoboszlai", "Maximizing open-play shot volume & direct set pieces."),
        ("VALUE ENABLERS", C_LAVENDER, "João Pedro • Calvert-Lewin", "Budget frontline unlocking premium midfield & defense.")
    ]
    for idx, (tag, col, title, desc) in enumerate(squad_cards):
        top_y = Inches(2.4) + idx * Inches(1.9)
        sc_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top_y, card_w4, Inches(1.7))
        sc_box.fill.solid()
        sc_box.fill.fore_color.rgb = RGBColor(28, 15, 60)
        sc_box.line.color.rgb = col
        
        stf = sc_box.text_frame
        stf.word_wrap = True
        stf.margin_left = Inches(0.18)
        stf.margin_top = Inches(0.15)
        
        sp1 = stf.paragraphs[0]
        sp1.text = tag
        sp1.font.size = Pt(11)
        sp1.font.bold = True
        sp1.font.color.rgb = col
        
        sp2 = stf.add_paragraph()
        sp2.text = title
        sp2.font.size = Pt(14)
        sp2.font.bold = True
        sp2.font.color.rgb = C_WHITE
        sp2.space_before = Pt(4)

        sp3 = stf.add_paragraph()
        sp3.text = desc
        sp3.font.size = Pt(11)
        sp3.font.color.rgb = C_SLATE
        sp3.space_before = Pt(4)

    # Right: Embedded Screenshot
    if os.path.exists(screenshot_path):
        s4.shapes.add_picture(screenshot_path, Inches(3.9), Inches(2.4), width=Inches(3.6))

    # Banner
    ban4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(8.3), Inches(7.0), Inches(0.7))
    ban4.fill.solid()
    ban4.fill.fore_color.rgb = RGBColor(28, 15, 60)
    ban4.line.color.rgb = C_ORCHID
    bp4 = ban4.text_frame.paragraphs[0]
    bp4.text = "£100.0M Cost  •  3-5-2 Formation  •  Max 3 per Club  •  5-Week Horizon"
    bp4.font.size = Pt(13)
    bp4.font.bold = True
    bp4.font.color.rgb = C_WHITE
    bp4.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 5: WHY RAW STATS LIE (Tactical Diagram)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    bg5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.0), Inches(10.0))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = C_SUNSET_DARK
    bg5.line.color.rgb = C_SUNSET_DARK

    add_header_footer(s5, 5)

    t_box5 = s5.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(7.0), Inches(1.2))
    tf5 = t_box5.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "WHY RAW STATS LIE"
    p5.font.size = Pt(26)
    p5.font.bold = True
    p5.font.color.rgb = C_WHITE

    p5_sub = tf5.add_paragraph()
    p5_sub.text = "Algorithms get tricked by shiny stats unless you enforce reality constraints."
    p5_sub.font.size = Pt(14)
    p5_sub.font.color.rgb = C_SLATE
    p5_sub.space_before = Pt(4)

    # Left: 3 Filter Cards
    card_w5 = Inches(3.4)
    filters = [
        ("CALIBRATION 01", C_TANGERINE, "Promoted Teams Reality Tax", "20% discount on lower-tier shot creation rates against Premier League defenses."),
        ("CALIBRATION 02", C_BLUE, "Bayes Sample-Size Shrinkage", "Blends 6-week short form with multi-season baselines to stop 1-game overreactions."),
        ("CALIBRATION 03", C_MINT, "Rotation & Fatigue Decay", "10% playing probability reduction for second matches in Double Gameweeks.")
    ]
    for idx, (tag, col, title, desc) in enumerate(filters):
        top_y = Inches(2.4) + idx * Inches(1.9)
        fc_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top_y, card_w5, Inches(1.7))
        fc_box.fill.solid()
        fc_box.fill.fore_color.rgb = RGBColor(35, 18, 12)
        fc_box.line.color.rgb = col
        
        ftf = fc_box.text_frame
        ftf.word_wrap = True
        ftf.margin_left = Inches(0.18)
        ftf.margin_top = Inches(0.15)
        
        fp1 = ftf.paragraphs[0]
        fp1.text = tag
        fp1.font.size = Pt(11)
        fp1.font.bold = True
        fp1.font.color.rgb = col
        
        fp2 = ftf.add_paragraph()
        fp2.text = title
        fp2.font.size = Pt(14)
        fp2.font.bold = True
        fp2.font.color.rgb = C_WHITE
        fp2.space_before = Pt(4)

        fp3 = ftf.add_paragraph()
        fp3.text = desc
        fp3.font.size = Pt(11)
        fp3.font.color.rgb = C_SLATE
        fp3.space_before = Pt(4)

    # Right: Opposition Matchup Box
    match_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(2.4), Inches(3.4), Inches(5.5))
    match_box.fill.solid()
    match_box.fill.fore_color.rgb = RGBColor(28, 14, 10)
    match_box.line.color.rgb = C_TANGERINE
    
    mtf5 = match_box.text_frame
    mtf5.word_wrap = True
    mtf5.margin_left = Inches(0.2)
    mtf5.margin_top = Inches(0.25)
    
    mp1 = mtf5.paragraphs[0]
    mp1.text = "TACTICAL OPPOSITION FILTER"
    mp1.font.size = Pt(13)
    mp1.font.bold = True
    mp1.font.color.rgb = C_WHITE

    items_m5 = [
        ("Premier League Defense Baseline", "Elite pressing & shot suppression floor."),
        ("The -20% Translation Tax", "Adjusts Championship shot volume to top-flight quality."),
        ("Clean Sheet Odds Ceiling", "Caps clean sheet probability at 24.6% max for promoted defenses.")
    ]
    for mt, md in items_m5:
        p_t = mtf5.add_paragraph()
        p_t.text = mt
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_AMBER
        p_t.space_before = Pt(18)
        
        p_d = mtf5.add_paragraph()
        p_d.text = md
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_SLATE
        p_d.space_before = Pt(3)

    # Banner
    ban5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(8.3), Inches(7.0), Inches(0.7))
    ban5.fill.solid()
    ban5.fill.fore_color.rgb = RGBColor(35, 18, 12)
    ban5.line.color.rgb = C_TANGERINE
    bp5 = ban5.text_frame.paragraphs[0]
    bp5.text = "Turn statistical noise and recency bias into pure predictive signal."
    bp5.font.size = Pt(13)
    bp5.font.bold = True
    bp5.font.color.rgb = C_WHITE
    bp5.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 6: THE CLOSING BENCHMARK & CTA (Text Only)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    bg6 = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.0), Inches(10.0))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = C_OBSIDIAN
    bg6.line.color.rgb = C_OBSIDIAN

    add_header_footer(s6, 6)

    t_box6 = s6.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(7.0), Inches(1.2))
    tf6 = t_box6.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "37 GAMEWEEKS. ZERO OVERRIDES."
    p6.font.size = Pt(28)
    p6.font.bold = True
    p6.font.color.rgb = C_WHITE

    p6_sub = tf6.add_paragraph()
    p6_sub.text = "Starting from dead last with 0 points. Handing the keys completely to the code."
    p6_sub.font.size = Pt(14)
    p6_sub.font.color.rgb = C_SLATE
    p6_sub.space_before = Pt(4)

    # 4 Rules 2x2 Grid
    r_w = Inches(3.35)
    r_h = Inches(1.6)
    rules = [
        ("RULE 01", "Model selects Starting XI, Captain & Vice Captain."),
        ("RULE 02", "Model executes all free transfers, hits, and chips."),
        ("RULE 03", "Zero human overrides, even when I hate the pick."),
        ("RULE 04", "Transparent weekly post-mortems every Tuesday.")
    ]
    positions = [
        (Inches(0.5), Inches(2.4)),
        (Inches(4.15), Inches(2.4)),
        (Inches(0.5), Inches(4.2)),
        (Inches(4.15), Inches(4.2))
    ]
    for i, (rx, ry) in enumerate(positions):
        rtag, rtxt = rules[i]
        rc = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, r_w, r_h)
        rc.fill.solid()
        rc.fill.fore_color.rgb = RGBColor(16, 28, 45)
        rc.line.color.rgb = C_MINT
        
        rtf = rc.text_frame
        rtf.word_wrap = True
        rtf.margin_left = Inches(0.18)
        rtf.margin_top = Inches(0.15)
        
        rp1 = rtf.paragraphs[0]
        rp1.text = rtag
        rp1.font.size = Pt(11)
        rp1.font.bold = True
        rp1.font.color.rgb = C_MINT
        
        rp2 = rtf.add_paragraph()
        rp2.text = rtxt
        rp2.font.size = Pt(13)
        rp2.font.bold = True
        rp2.font.color.rgb = C_WHITE
        rp2.space_before = Pt(4)

    # Grand CTA Box
    cta_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(7.0), Inches(2.1))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = RGBColor(10, 36, 28)
    cta_box.line.color.rgb = C_MINT
    
    ctf6 = cta_box.text_frame
    ctf6.word_wrap = True
    ctf6.margin_left = Inches(0.25)
    ctf6.margin_top = Inches(0.2)
    
    cp1 = ctf6.paragraphs[0]
    cp1.text = "THE BENCHMARK QUESTION"
    cp1.font.size = Pt(17)
    cp1.font.bold = True
    cp1.font.color.rgb = C_WHITE
    
    cp2 = ctf6.add_paragraph()
    cp2.text = "Where does an automated algorithm finish by May after spotting 11 million managers a 1-week head start?"
    cp2.font.size = Pt(13)
    cp2.font.color.rgb = C_WHITE
    cp2.space_before = Pt(6)

    cp3 = ctf6.add_paragraph()
    cp3.text = "🔥 TOP 100K    •    📈 TOP 1M    •    💥 TOTAL DISASTER"
    cp3.font.size = Pt(13)
    cp3.font.bold = True
    cp3.font.color.rgb = C_AMBER
    cp3.space_before = Pt(10)

    # Banner
    ban6 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(8.3), Inches(7.0), Inches(0.7))
    ban6.fill.solid()
    ban6.fill.fore_color.rgb = RGBColor(10, 36, 28)
    ban6.line.color.rgb = C_MINT
    bp6 = ban6.text_frame.paragraphs[0]
    bp6.text = "Drop your prediction in the comments. Weekly updates begin next Tuesday."
    bp6.font.size = Pt(13)
    bp6.font.bold = True
    bp6.font.color.rgb = C_WHITE
    bp6.alignment = PP_ALIGN.CENTER

    out_path = os.path.abspath("carousels/fpl-algorithm-experiment/fpl-algorithm-experiment.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully to: {out_path}")

if __name__ == "__main__":
    create_fpl_pptx()
